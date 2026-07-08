from pptx import Presentation
from pptx.util import Inches, Pt
import os

def add_title_slide(prs, title, subtitle, team):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle + "\n\n" + team
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Greet jury; state project name and goal."

def add_bullet_slide(prs, title, bullets, notes_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 1
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

def build_presentation(path):
    prs = Presentation()

    add_title_slide(prs,
                    "GhostHunter SOC",
                    "Autonomous SOC for faster detection & response",
                    "Team: GhostHunter — Demo for Jury | TCS - TATA Consultancy Services")

    add_bullet_slide(prs, "Problem",
                     ["High alert volumes and long investigation times",
                      "Manual triage is slow and error-prone",
                      "Enterprises need scalable, automated SOC workflows"],
                     "Speak: Describe the pain — noisy alerts, long MTTR. Mention typical daily alert volumes and analyst load; set expectation for demo showing automation benefits.")

    add_bullet_slide(prs, "Solution Summary",
                     ["Autonomous pipeline: ingestion → detection → correlation → response",
                      "LLM-assisted agents for investigation and reporting",
                      "Graph-backed context (Neo4j) + semantic store (Chroma)"],
                     "Speak: Outline the automated flow and how each component reduces manual work. Emphasize novelty: combining graph context with RAG and agents to automate investigation.")

    add_bullet_slide(prs, "Architecture & Tech",
                     ["Data ingestion: Kafka-based pipeline",
                      "Detection: rule + ML models",
                      "Knowledge: Neo4j attack graph, Chroma RAG store",
                      "Agents: automated response, reporting, correlation"],
                     "Speak: Show architecture diagram. Call out technologies and why chosen (scalability, semantic search, graph relations). Mention data flow and latency expectations.")

    add_bullet_slide(prs, "Demo Preview",
                     ["Live demo: detection → investigation → response",
                      "Fallback: short recorded clip (30s)",
                      "Screenshots included in backup slides"],
                     "Speak: Briefly describe the demo steps. If live, we'll trigger an alert and show the agent investigation and generated report. If recorded, play the clip and narrate key moments.")

    add_bullet_slide(prs, "Results & ROI",
                     ["Reduced mean time to respond (example: 60%)",
                      "Fewer analyst hours per incident",
                      "Improved detection precision and context for triage"],
                     "Speak: Present conservative estimates: percentage MTTR reduction, analyst hours saved, and potential cost savings. Tie to business impact and compliance benefits.")

    add_bullet_slide(prs, "Roadmap & Risks",
                     ["Next: field pilot, integration with SIEMs, hardened ops",
                      "Risks: model drift, integration efforts; mitigations planned"],
                     "Speak: Show 3-month roadmap and key milestones. Acknowledge risks and mitigation (monitoring, model retraining, phased rollouts).")

    add_bullet_slide(prs, "Ask & Contact",
                     ["Request: pilots, feedback, potential partners/funding",
                      "Contact: ghosthunter-team@example.com"],
                     "Speak: State specific asks (pilot customers, integration partners, seed funding). Provide contact info and invite follow-up demos.")

    # Add a concise handout slide summarizing value prop (for printing or PDF handout)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "One-page Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = "Problem: Alert overload and slow investigations"
    p = body.add_paragraph()
    p.text = "Solution: Autonomous SOC combining graph context, RAG, and agent automation"
    p.level = 1
    p2 = body.add_paragraph()
    p2.text = "Ask: Pilots & partners; contact ghosthunter-team@example.com"
    p2.level = 1
    prs.slides[-1].notes_slide.notes_text_frame.text = "Handout: Give this one-page summary as a takeaway for judges."

    # Add screenshots from images folder as backup slides
    def add_image_slides(prs, images_dir='images'):
        if not os.path.isdir(images_dir):
            return
        imgs = [f for f in os.listdir(images_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        imgs.sort()
        per_slide = 2
        for i in range(0, len(imgs), per_slide):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            chunk = imgs[i:i+per_slide]
            for j, img in enumerate(chunk):
                img_path = os.path.join(images_dir, img)
                left = Inches(0.5) if j == 0 else Inches(5)
                top = Inches(1)
                width = Inches(4.0)
                try:
                    slide.shapes.add_picture(img_path, left, top, width=width)
                except Exception:
                    # skip images that can't be inserted
                    pass

    add_image_slides(prs, images_dir='images')

    # Save
    prs.save(path)

if __name__ == '__main__':
    out_path = 'presentation/GhostHunter_Jury_Presentation.pptx'
    build_presentation(out_path)
    print(f'Created presentation: {out_path}')
